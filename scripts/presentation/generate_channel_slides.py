"""
Generate PowerPoint slides for EEG channel reduction experiments.

Usage:
    uv run python scripts/presentation/generate_channel_slides.py
    uv run python scripts/presentation/generate_channel_slides.py --output custom_path.pptx
"""

import argparse
import tempfile
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================================
# Constants
# ============================================================================

PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors matching project constants (src/config/constants.py)
COLOR_EEGNET = RGBColor(0x2E, 0x86, 0xAB)
COLOR_CBRAMOD = RGBColor(0xE9, 0x4F, 0x37)
COLOR_HEADER = RGBColor(0x2C, 0x3E, 0x50)
COLOR_ACCENT = RGBColor(0x27, 0xAE, 0x60)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_MID_GRAY = RGBColor(0x99, 0x99, 0x99)
COLOR_WARNING = RGBColor(0xE6, 0x7E, 0x22)

FONT_TITLE = 'Microsoft YaHei'
FONT_BODY = 'Microsoft YaHei'
FONT_MONO = 'Consolas'


# ============================================================================
# Helper functions
# ============================================================================

def set_slide_bg(slide, color=COLOR_WHITE):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a colored title bar at the top of a content slide."""
    # Title background bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.1),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_HEADER
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.15),
        Inches(12), Inches(0.55),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = FONT_TITLE

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.7),
            Inches(12), Inches(0.35),
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p2.font.name = FONT_BODY


def add_text_block(slide, left, top, width, height, text_lines, font_size=18,
                   bold=False, color=COLOR_DARK_TEXT, line_spacing=1.5):
    """Add a multi-line text block."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(font_size * (line_spacing - 1))


def add_bullet_points(slide, left, top, width, height, items, font_size=18,
                      color=COLOR_DARK_TEXT):
    """Add bullet point list."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle nested items (tuple: (text, indent_level))
        if isinstance(item, tuple):
            text, level = item
            p.text = text
            p.level = level
        else:
            p.text = f"  {item}"
            p.level = 0

        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(6)


def add_table(slide, left, top, width, height, headers, rows,
              header_color=COLOR_HEADER, highlight_rows=None):
    """Add a styled table to a slide."""
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    table = table_shape.table

    # Set column widths proportionally
    col_width = Inches(width / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_width

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLOR_WHITE
            paragraph.font.name = FONT_BODY
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        is_highlight = highlight_rows and i in highlight_rows
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if is_highlight:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = COLOR_DARK_TEXT
                paragraph.font.name = FONT_BODY
                paragraph.alignment = PP_ALIGN.CENTER
                if is_highlight:
                    paragraph.font.bold = True
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """Add an image if it exists, otherwise add placeholder text."""
    img = Path(img_path)
    if img.exists():
        kwargs = {}
        if width:
            kwargs['width'] = Inches(width)
        if height:
            kwargs['height'] = Inches(height)
        slide.shapes.add_picture(
            str(img), Inches(left), Inches(top), **kwargs
        )
        return True
    else:
        add_text_block(slide, left, top, 6, 1,
                       [f"[Image not found: {img.name}]"],
                       font_size=12, color=COLOR_MID_GRAY)
        return False


def generate_degradation_curve(output_path: str):
    """Generate the channel count vs accuracy degradation curve."""
    plt.rcParams['font.family'] = ['Microsoft YaHei', 'sans-serif']

    channels = [8, 32, 128]

    # CBraMod best results (transfer where available, cross-subject for 128ch)
    binary_transfer = [72.92, 88.90, 90.27]
    ternary_transfer = [57.26, 72.68, 75.42]
    binary_cross = [68.33, 88.10, 90.27]
    ternary_cross = [52.00, 70.79, 75.42]

    fig, ax = plt.subplots(figsize=(10, 6))

    # Transfer lines (solid)
    ax.plot(channels, binary_transfer, 'o-', color='#E94F37', linewidth=2.5,
            markersize=12, label='Binary Transfer', zorder=5)
    ax.plot(channels, ternary_transfer, 's-', color='#2E86AB', linewidth=2.5,
            markersize=12, label='Ternary Transfer', zorder=5)

    # Cross-subject lines (dashed)
    ax.plot(channels, binary_cross, 'o--', color='#E94F37', linewidth=1.5,
            markersize=8, alpha=0.5, label='Binary Cross-Subject')
    ax.plot(channels, ternary_cross, 's--', color='#2E86AB', linewidth=1.5,
            markersize=8, alpha=0.5, label='Ternary Cross-Subject')

    # Sweet spot annotation
    ax.axvspan(24, 40, alpha=0.10, color='#27AE60', zorder=0)
    ax.annotate('Sweet Spot\n(32ch)', xy=(32, 92), fontsize=13,
                color='#27AE60', fontweight='bold', ha='center')

    # Performance cliff annotation
    ax.annotate('', xy=(15, 62), xytext=(15, 85),
                arrowprops=dict(arrowstyle='<->', color='#E94F37', lw=2))
    ax.text(11.5, 73, 'Performance\nCliff', fontsize=11, color='#E94F37',
            ha='center', fontweight='bold')

    # Data labels
    for ch, ba, ta in zip(channels, binary_transfer, ternary_transfer):
        ax.annotate(f'{ba:.1f}%', xy=(ch, ba), textcoords='offset points',
                    xytext=(12, 8), fontsize=11, color='#E94F37', fontweight='bold')
        ax.annotate(f'{ta:.1f}%', xy=(ch, ta), textcoords='offset points',
                    xytext=(12, -15), fontsize=11, color='#2E86AB', fontweight='bold')

    # Chance level
    ax.axhline(y=50, color='gray', linestyle=':', alpha=0.5, linewidth=1)
    ax.text(130, 50.8, 'Chance (Binary)', fontsize=9, color='gray', alpha=0.7)
    ax.axhline(y=33.3, color='gray', linestyle=':', alpha=0.5, linewidth=1)
    ax.text(130, 34.1, 'Chance (Ternary)', fontsize=9, color='gray', alpha=0.7)

    ax.set_xlabel('Number of EEG Channels', fontsize=14)
    ax.set_ylabel('Mean Test Accuracy (%)', fontsize=14)
    ax.set_xscale('log', base=2)
    ax.set_xticks(channels)
    ax.set_xticklabels([str(c) for c in channels])
    ax.legend(fontsize=12, loc='lower right')
    ax.grid(True, alpha=0.3)
    ax.set_ylim(25, 100)
    ax.set_xlim(6, 180)
    ax.set_title('CBraMod: Channel Count vs Decoding Accuracy',
                 fontsize=16, fontweight='bold', pad=15)

    fig.tight_layout()
    fig.savefig(output_path, dpi=200, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    return output_path


# ============================================================================
# Slide builders
# ============================================================================

def add_slide_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Accent bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(2.5),
        SLIDE_WIDTH, Inches(3.0),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_HEADER
    shape.line.fill.background()

    # Title
    add_text_block(slide, 0.8, 2.7, 12, 1.2,
                   ['EEG 通道缩减实验'],
                   font_size=40, bold=True, color=COLOR_WHITE)
    add_text_block(slide, 0.8, 3.7, 12, 0.8,
                   ['128 → 32 → 8 通道: CBraMod vs EEGNet 性能权衡分析'],
                   font_size=22, color=RGBColor(0xBB, 0xBB, 0xBB))
    add_text_block(slide, 0.8, 4.6, 12, 0.5,
                   ['Motor Imagery — 单指级运动想象解码  |  21 被试  |  BioSemi 128 EEG'],
                   font_size=16, color=RGBColor(0x99, 0x99, 0x99))

    # Date
    add_text_block(slide, 0.8, 6.5, 12, 0.5,
                   ['2026-02'],
                   font_size=14, color=COLOR_MID_GRAY)


def add_slide_motivation(prs):
    """Slide 2: Research motivation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, '研究动机', 'Why reduce EEG channels?')

    items = [
        '高密度 EEG 设备 (128ch) 昂贵、准备时间长、不适合日常使用',
        '商用 BCI 设备通常仅 8-32 通道',
        '核心问题: 在保持解码性能的前提下，最少需要多少通道？',
        '附加问题: 通道选择策略 — 手工选择 vs 数据驱动？',
    ]
    add_bullet_points(slide, 0.8, 1.5, 11.5, 2.5, items, font_size=22)

    # Diagram: 128 -> 32 -> 8
    for i, (ch, x) in enumerate([(128, 2.0), (32, 5.5), (8, 9.0)]):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(4.5),
            Inches(2.2), Inches(1.5),
        )
        shape.fill.solid()
        colors = [COLOR_CBRAMOD, COLOR_ACCENT, COLOR_EEGNET]
        shape.fill.fore_color.rgb = colors[i]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f'{ch} 通道'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        labels = ['研究级', '商用级', '消费级']
        p2.text = labels[i]
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        p2.font.name = FONT_BODY
        p2.alignment = PP_ALIGN.CENTER

    # Arrows
    for x in [4.3, 7.8]:
        add_text_block(slide, x, 4.9, 1.2, 0.8, ['→'],
                       font_size=36, bold=True, color=COLOR_HEADER)


def add_slide_overview(prs):
    """Slide 3: Experiment overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, '实验设计概览')

    # Left: Pipeline
    add_text_block(slide, 0.8, 1.4, 5, 0.5,
                   ['实验管线 (Pipeline)'], font_size=20, bold=True)
    pipeline_items = [
        '1. Cross-Subject 预训练 (21 被试, Leave-One-Out 评估)',
        '2. Transfer Learning (预训练 → 个体微调)',
        '3. 6 种 32ch 配置对比 → 最优配置全量实验',
    ]
    add_bullet_points(slide, 0.8, 1.9, 6, 2, pipeline_items, font_size=16)

    # Right: Models & Tasks
    add_table(slide, 7.2, 1.4, 5.5, 1.8,
              ['', 'CBraMod', 'EEGNet'],
              [
                  ['类型', '预训练基座模型', '传统 CNN'],
                  ['参数量', '~4.0M', '~2.5K'],
                  ['采样率', '200 Hz', '100 Hz'],
                  ['滤波', '0.3-75 Hz', '4-40 Hz'],
              ])

    # Tasks
    add_table(slide, 0.8, 4.2, 12, 1.5,
              ['任务', '类别', '描述', '随机基线'],
              [
                  ['Binary', '拇指 vs 小指', '2 类运动想象', '50%'],
                  ['Ternary', '拇指 vs 中指 vs 小指', '3 类运动想象', '33.3%'],
              ])

    # Channel configs summary
    add_text_block(slide, 0.8, 6.0, 12, 0.5,
                   ['通道配置: 2 种手工选择 (motor_cortex, commercial) + '
                    '4 种数据驱动 (FDR, CSP, Attention, Band Power)'],
                   font_size=15, color=COLOR_MID_GRAY)


def add_slide_baseline(prs):
    """Slide 4: 128ch baseline results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, '128 通道基线结果', 'Full BioSemi 128 — Cross-Subject CBraMod')

    add_table(slide, 0.8, 1.5, 11.5, 1.2,
              ['任务', 'Cross-Subject (Mean±Std)', 'Median', 'Min', 'Max', '被试数'],
              [
                  ['Binary (拇指 vs 小指)', '90.27 ± 8.88%', '93.12%', '66.25%', '99.38%', '21'],
                  ['Ternary (3 指)', '75.42 ± 13.27%', '77.50%', '43.75%', '93.75%', '21'],
              ],
              highlight_rows=[0])

    add_bullet_points(slide, 0.8, 3.5, 11, 3, [
        'CBraMod 基座模型在全通道下表现优异: Binary 90.27%',
        'Ternary 相比 Binary 下降 ~15pp — 三指解码显著更难',
        '被试间差异大 (S20: 66.25%, S19: 99.38%) — 反映 EEG 个体差异',
        '此为后续通道缩减实验的对照基线',
    ], font_size=18)


def add_slide_strategies(prs):
    """Slide 5: 6 channel selection strategies."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, '32ch 通道选择策略', '6 configurations compared')

    add_table(slide, 0.4, 1.4, 12.5, 3.5,
              ['配置名', '类型', '方法', '特点'],
              [
                  ['motor_cortex', 'Hand-picked',
                   'C3/Cz/C4 + SMA + premotor 密集覆盖',
                   '基于运动皮层先验知识'],
                  ['commercial', 'Hand-picked',
                   '标准 10-20 布局 (Fp1, F3, Fz, C3, Cz...)',
                   '模拟商用 32ch EEG 帽'],
                  ['FDR', 'Data-driven',
                   'Fisher Discriminant Ratio 通道排序',
                   '纯统计,无需模型,极快'],
                  ['CSP', 'Data-driven',
                   'Common Spatial Patterns 权重排序',
                   '考虑空间滤波器权重'],
                  ['Attention', 'Data-driven',
                   'EEGNet 权重 + CBraMod 梯度',
                   '结合两模型的通道注意力'],
                  ['Band Power', 'Data-driven',
                   'Mu(8-13Hz) + Beta(13-30Hz) ANOVA',
                   '基于运动相关频段能量'],
              ])

    add_text_block(slide, 0.8, 5.8, 11, 1,
                   ['FDR 公式: FDR_ch = (μ₁ - μ₂)² / (σ₁² + σ₂²)  — '
                    '类间距离 / 类内散度'],
                   font_size=15, color=COLOR_MID_GRAY)


def add_slide_32ch_comparison_chart(prs):
    """Slide 6: 32ch 6-config comparison chart (embed existing image)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, '32ch 6 配置对比', 'Cross-Subject Binary — CBraMod + EEGNet')

    img_path = PROJECT_ROOT / 'results' / '32_channel' / \
        '20260222_1324_32ch_config_comparison_imagery_binary.png'
    add_image_safe(slide, img_path, 0.5, 1.3, width=12.3)


def add_slide_32ch_ranking(prs):
    """Slide 7: 32ch ranking table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, '32ch 配置排名', 'Binary Cross-Subject — 综合排名')

    add_table(slide, 0.8, 1.5, 11.5, 3.5,
              ['排名', '配置', '类型', 'CBraMod', 'EEGNet', '综合均值'],
              [
                  ['1', 'Attention', 'Data-driven', '87.02%', '70.42%', '78.72%'],
                  ['2', 'FDR', 'Data-driven', '88.10%', '67.53%', '77.81%'],
                  ['3', 'Band Power', 'Data-driven', '85.51%', '67.17%', '76.34%'],
                  ['4', 'CSP', 'Data-driven', '85.54%', '66.52%', '76.03%'],
                  ['5', 'Commercial', 'Hand-picked', '86.31%', '64.40%', '75.36%'],
                  ['6', 'Motor Cortex', 'Hand-picked', '82.02%', '63.12%', '72.57%'],
              ],
              highlight_rows=[0, 1])

    add_bullet_points(slide, 0.8, 5.5, 11, 1.5, [
        '数据驱动方法全面优于手工选择 — 4 个 data-driven 均排在 hand-picked 之前',
        'CBraMod 最优: FDR (88.10%), 仅比 128ch 基线低 2.17%',
        'FDR 为最简单的统计方法，无需模型，但在 CBraMod 上排名第一',
    ], font_size=16)


def add_slide_32ch_binary_pipeline(prs):
    """Slide 8: 32ch FDR binary full pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, '32ch FDR: Binary 全管线',
                  'Cross-Subject → Transfer Learning (CBraMod)')

    img_path = PROJECT_ROOT / 'results' / '32_channel' / \
        '20260221_0445_transfer_combined_imagery_binary.png'
    add_image_safe(slide, img_path, 0.3, 1.3, width=7.5)

    # Side summary
    add_table(slide, 8.2, 1.5, 4.5, 1.5,
              ['方法', 'Mean Accuracy', 'vs 128ch'],
              [
                  ['Cross-Subject', '88.10%', '-2.17%'],
                  ['Transfer', '88.90%', '-1.37%'],
                  ['128ch 基线', '90.27%', '—'],
              ],
              highlight_rows=[1])

    add_bullet_points(slide, 8.2, 3.5, 4.5, 3.5, [
        '75% 通道削减 → 仅 2.17% 准确率损失',
        'Transfer 微调额外恢复 +0.80%',
        '低表现被试受益最大 (S05: +9.38%)',
        '训练时间减半 (~44min vs ~85min)',
    ], font_size=15)


def add_slide_32ch_ternary_pipeline(prs):
    """Slide 9: 32ch FDR ternary full pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, '32ch FDR: Ternary 全管线',
                  'Cross-Subject → Transfer Learning (CBraMod, 3-class)')

    img_path = PROJECT_ROOT / 'results' / '32_channel' / \
        '20260221_1042_transfer_combined_imagery_ternary.png'
    add_image_safe(slide, img_path, 0.3, 1.3, width=7.5)

    add_table(slide, 8.2, 1.5, 4.5, 1.5,
              ['方法', 'Mean Accuracy', 'vs 128ch'],
              [
                  ['Cross-Subject', '70.79%', '-4.62%'],
                  ['Transfer', '72.68%', '-2.74%'],
                  ['128ch 基线', '75.42%', '—'],
              ],
              highlight_rows=[1])

    add_bullet_points(slide, 8.2, 3.5, 4.5, 3.5, [
        'Ternary 通道损失大于 Binary (-4.62% vs -2.17%)',
        '三分类对空间分辨率更敏感',
        'Transfer 恢复约 41% 通道削减损失',
        'S05 ternary transfer 提升 +12.50%',
    ], font_size=15)


def add_slide_8ch_cliff(prs):
    """Slide 10: 8ch performance cliff."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, '8ch: 性能断崖', '32→8 通道的急剧下降 (CBraMod FDR)')

    # Binary image
    img_path = PROJECT_ROOT / 'results' / '8_channel' / \
        '20260221_1319_transfer_combined_imagery_binary.png'
    add_image_safe(slide, img_path, 0.2, 1.3, width=6.3)

    # Ternary image
    img_path2 = PROJECT_ROOT / 'results' / '8_channel' / \
        '20260221_1547_transfer_combined_imagery_ternary.png'
    add_image_safe(slide, img_path2, 6.7, 1.3, width=6.3)

    # Summary at bottom
    add_table(slide, 1.5, 5.5, 10, 1.3,
              ['任务', 'Cross-Subject', 'Transfer', 'Transfer Δ', 'vs 128ch'],
              [
                  ['Binary', '68.33%', '72.92%', '+4.59%', '-17.35%'],
                  ['Ternary', '52.00%', '57.26%', '+5.26%', '-18.16%'],
              ])


def add_slide_degradation_curve(prs, chart_path):
    """Slide 11: Channel degradation curve (matplotlib generated)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, '通道退化曲线',
                  'Channel Count vs Accuracy — CBraMod (FDR)')

    add_image_safe(slide, chart_path, 0.8, 1.3, width=8)

    # Key insight box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.2), Inches(1.5),
        Inches(3.8), Inches(3.0),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    shape.line.color.rgb = COLOR_ACCENT

    add_text_block(slide, 9.4, 1.7, 3.4, 0.5,
                   ['核心发现'], font_size=18, bold=True, color=COLOR_ACCENT)
    add_bullet_points(slide, 9.4, 2.3, 3.4, 2, [
        '128→32ch: 仅 -1.4% (Binary Transfer)',
        '32→8ch: -16.0% — 性能断崖',
        '32ch 是最优权衡点',
        '低通道数下 Transfer 效果更显著',
    ], font_size=14, color=COLOR_DARK_TEXT)


def add_slide_transfer_scaling(prs):
    """Slide 12: Transfer learning effectiveness scales with channel reduction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, 'Transfer 微调效果随通道数变化',
                  'Fewer channels → Greater benefit from individual adaptation')

    add_table(slide, 0.8, 1.5, 11.5, 2.0,
              ['通道数', 'Binary Cross', 'Binary Transfer', 'Δ',
               'Ternary Cross', 'Ternary Transfer', 'Δ'],
              [
                  ['128ch', '90.27%', '—', '—', '75.42%', '—', '—'],
                  ['32ch FDR', '88.10%', '88.90%', '+0.80%',
                   '70.79%', '72.68%', '+1.89%'],
                  ['8ch FDR', '68.33%', '72.92%', '+4.59%',
                   '52.00%', '57.26%', '+5.26%'],
              ],
              highlight_rows=[2])

    add_bullet_points(slide, 0.8, 4.2, 11, 3, [
        'Transfer 微调增益随通道减少而增大: '
        'Binary +0.80% (32ch) → +4.59% (8ch)',
        '通道越少, 跨被试模型泛化能力越弱, 个体适配价值越高',
        '极端案例: S05 在 8ch binary 从 45.62% 跃升至 70.62% (+25.00%)',
        'CBraMod vs EEGNet 差距也随通道减少扩大: '
        '32ch gap 20.57pp (88.10% vs 67.53%)',
    ], font_size=17)


def add_slide_caveats(prs):
    """Slide 13: Caveats and limitations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, '注意事项与局限性', 'Caveats & Limitations')

    add_text_block(slide, 0.8, 1.4, 11.5, 0.5,
                   ['1. 超参数并非完全控制变量'], font_size=20, bold=True,
                   color=COLOR_WARNING)
    add_bullet_points(slide, 1.0, 1.9, 11, 1.5, [
        'CBraMod 在低通道数下使用了更强的正则化 '
        '(8ch: dropout 0.30, weight_decay 0.10 vs 128ch: 0.15, 0.06)',
        'EEGNet 训练超参数未随通道数变化',
        '因此性能差异并非纯粹由通道数决定 — 包含超参数调优的影响',
    ], font_size=15)

    add_text_block(slide, 0.8, 3.4, 11.5, 0.5,
                   ['2. "Commercial" 配置不等于真实商用设备'],
                   font_size=20, bold=True, color=COLOR_WARNING)
    add_bullet_points(slide, 1.0, 3.9, 11, 1.5, [
        '数据来源仍是 BioSemi 128 研究级设备 (湿电极, 高 SNR)',
        '真实商用设备在传感器类型 (干/湿)、阻抗、前端处理上有显著差异',
        '本实验的 commercial 配置结果不能直接推广到实际商用 EEG 帽',
    ], font_size=15)

    add_text_block(slide, 0.8, 5.4, 11.5, 0.5,
                   ['3. 其他局限'], font_size=20, bold=True, color=COLOR_WARNING)
    add_bullet_points(slide, 1.0, 5.9, 11, 1.2, [
        '单一数据集 (21 被试, Motor Imagery) — 结论待跨数据集验证',
        'FDR 数据驱动选择使用了相同被试数据 — 存在潜在过拟合风险',
        '仅测试 Motor Imagery 范式, Motor Execution 待实验',
    ], font_size=15)


def add_slide_conclusions(prs):
    """Slide 14: Conclusions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, '结论与展望')

    add_text_block(slide, 0.8, 1.4, 11.5, 0.5,
                   ['核心发现'], font_size=22, bold=True)
    conclusions = [
        '32ch 是性能/硬件的最优权衡点 — Binary 仅损失 2.17%, '
        'Ternary 损失 4.62%',
        '数据驱动通道选择全面优于手工选择 (FDR 最简单但 CBraMod 最优)',
        'CBraMod 基座模型在低通道数下优势更显著 — '
        '预训练知识弥补信息损失',
        'Transfer Learning 在低通道数下效果更大 — '
        '个体适配补偿泛化不足',
        '8ch 存在性能断崖 (-20pp) — '
        '极低通道数消费级 BCI 仍需突破',
    ]
    add_bullet_points(slide, 0.8, 2.0, 11, 3, conclusions, font_size=18)

    add_text_block(slide, 0.8, 5.2, 11.5, 0.5,
                   ['未来工作'], font_size=22, bold=True)
    future = [
        'Motor Execution 范式验证',
        '跨数据集泛化 (不同 EEG 设备/被试群体)',
        '32ch commercial 配置全量实验 (含 Transfer)',
        '真实商用 EEG 设备在线 BCI 测试',
    ]
    add_bullet_points(slide, 0.8, 5.8, 11, 1.5, future, font_size=16,
                      color=COLOR_MID_GRAY)


# ============================================================================
# Main
# ============================================================================

def main():
    parser = argparse.ArgumentParser(
        description='Generate channel reduction experiment slides (.pptx)')
    parser.add_argument('--output', '-o', type=str,
                        default=str(PROJECT_ROOT / 'results' / 'presentation' /
                                    'channel_reduction_slides.pptx'),
                        help='Output .pptx file path')
    args = parser.parse_args()

    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    print("Generating slides...")

    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Generate matplotlib chart
    chart_path = output_path.parent / 'degradation_curve.png'
    print(f"  Generating degradation curve → {chart_path}")
    generate_degradation_curve(str(chart_path))

    # Build slides
    print("  Building slides...")
    add_slide_title(prs)           # 1
    add_slide_motivation(prs)      # 2
    add_slide_overview(prs)        # 3
    add_slide_baseline(prs)        # 4
    add_slide_strategies(prs)      # 5
    add_slide_32ch_comparison_chart(prs)  # 6
    add_slide_32ch_ranking(prs)    # 7
    add_slide_32ch_binary_pipeline(prs)   # 8
    add_slide_32ch_ternary_pipeline(prs)  # 9
    add_slide_8ch_cliff(prs)       # 10
    add_slide_degradation_curve(prs, str(chart_path))  # 11
    add_slide_transfer_scaling(prs)  # 12
    add_slide_caveats(prs)         # 13
    add_slide_conclusions(prs)     # 14

    # Save
    prs.save(str(output_path))
    print(f"\nSlides saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
